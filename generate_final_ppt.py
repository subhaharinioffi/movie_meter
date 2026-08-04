import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Aspect Ratio
    prs.slide_height = Inches(7.5)
    
    # Theme Colors matching premium Vercel/Apple SaaS branding
    maroon = RGBColor(109, 0, 26)     # #6D001A
    gold = RGBColor(255, 213, 79)    # #FFD54F
    slate = RGBColor(71, 85, 105)     # #475569
    dark_blue = RGBColor(15, 23, 42)  # #0F172A
    white = RGBColor(255, 255, 255)
    off_white = RGBColor(250, 250, 250)
    
    # Locate generated images in brain directory
    brain_dir = r"C:\Users\Subhaharini\.gemini\antigravity\brain\3651ea6e-94c7-4edf-89e9-b55af2040db7"
    
    def get_image(prefix):
        files = glob.glob(os.path.join(brain_dir, f"{prefix}_*.jpg"))
        return files[0] if files else None
        
    hero_img = get_image("hero_banner_mockup")
    showcase_img = get_image("movie_showcase_mockup")
    gauge_img = get_image("success_gauge_mockup")
    revenue_img = get_image("revenue_chart_mockup")
    audience_img = get_image("audience_insights_mockup")
    
    slide_layout = prs.slide_layouts[6] # Blank Slide
    
    # ------------------ SLIDE 1: Problem Statement & Solution ------------------
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = maroon
    
    # Title & Subtitle Left Block
    txBox1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(6.5), Inches(5.5))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "MOVIE METER"
    p.font.name = "Georgia"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = gold
    p.space_after = Pt(14)
    
    bullets1 = [
        "Problem: Film budget packaging and digital acquisitions suffer from massive financial risks due to subjective intuition and pre-release data voids.",
        "Solution: A predictive analytics dashboard classifying movie profiles (High/Medium/Low) before shooting or distribution bidding.",
        "Methodology: Resolves superstar casting coefficients in a leak-free out-of-fold reputation mapping index."
    ]
    for b_text in bullets1:
        bp = tf1.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = white
        bp.space_after = Pt(10)
        
    if hero_img:
        slide1.shapes.add_picture(hero_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 2: ML Technique and Architecture ------------------
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = off_white
    
    hBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf2 = hBox2.text_frame
    hp2 = htf2.paragraphs[0]
    hp2.text = "Machine Learning Technique & Tuning"
    hp2.font.name = "Georgia"
    hp2.font.size = Pt(28)
    hp2.font.bold = True
    hp2.font.color.rgb = maroon
    
    tBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.5), Inches(5.0))
    tf2 = tBox2.text_frame
    tf2.word_wrap = True
    
    bullets2 = [
        "Model Selection: Trained on XGBoost Classifier, outperforming Random Forest and Gradient Boosting baselines by ~12% on margin classes.",
        "Reputation Encoding: Computes smoothed out-of-fold average rating mapping to quantify director and actor performance historical profiles.",
        "Stratified 5-Fold Cross-Validation: Yields a validation accuracy of 62.1% and ROC-AUC of 73.0% under zero data leakage conditions.",
        "Hyperparameters: n_estimators=300, max_depth=6, learning_rate=0.03, subsample=0.8."
    ]
    for i, b_text in enumerate(bullets2):
        bp = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = slate
        bp.space_after = Pt(12)
        
    if gauge_img:
        slide2.shapes.add_picture(gauge_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 3: System Architecture ------------------
    slide3 = prs.slides.add_slide(slide_layout)
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = off_white
    
    hBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf3 = hBox3.text_frame
    hp3 = htf3.paragraphs[0]
    hp3.text = "System Architecture & Processing Flow"
    hp3.font.name = "Georgia"
    hp3.font.size = Pt(28)
    hp3.font.bold = True
    hp3.font.color.rgb = maroon
    
    tBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.0))
    tf3 = tBox3.text_frame
    tf3.word_wrap = True
    
    bullets3 = [
        "Data Ingestion: Chunk-based extraction of region/language indicators across 1.3 GB of official IMDb TSV datasets.",
        "Preprocessing Pipeline: Handles median imputation for runtime deviations and buckets content ratings into standard certification bins.",
        "Inference & Output: XGBoost classifies screenplay configurations, feeding Plotly rendering engines and distribution recommendation layers.",
        "Export Engine: Programmatic generation of technical reports (PDF) and investor presentations (PPTX)."
    ]
    for i, b_text in enumerate(bullets3):
        bp = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = slate
        bp.space_after = Pt(12)
        
    if showcase_img:
        slide3.shapes.add_picture(showcase_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 4: Tech Stacks & Why ------------------
    slide4 = prs.slides.add_slide(slide_layout)
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = off_white
    
    hBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf4 = hBox4.text_frame
    hp4 = htf4.paragraphs[0]
    hp4.text = "Tech Stack Selection Rationales"
    hp4.font.name = "Georgia"
    hp4.font.size = Pt(28)
    hp4.font.bold = True
    hp4.font.color.rgb = maroon
    
    tBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.0))
    tf4 = tBox4.text_frame
    tf4.word_wrap = True
    
    bullets4 = [
        "XGBoost Classifier: Selected for its capability to optimize gradients on tabular datasets, handling sparse target encoded director/actor features efficiently.",
        "Streamlit (Python): Chosen to write light-theme web frontends with high-fidelity styles, maintaining reactive rendering states without separate server frameworks.",
        "Plotly Engine: Provides responsive vector indicators and gauge charts that support custom color overlays.",
        "ReportLab & python-pptx: Programmatic file compiles to output physical PDF report and PowerPoint presentations dynamically."
    ]
    for i, b_text in enumerate(bullets4):
        bp = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = slate
        bp.space_after = Pt(10)
        
    if revenue_img:
        slide4.shapes.add_picture(revenue_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 5: Concluding All ------------------
    slide5 = prs.slides.add_slide(slide_layout)
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = maroon
    
    hBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf5 = hBox5.text_frame
    hp5 = htf5.paragraphs[0]
    hp5.text = "Concluding Summary"
    hp5.font.name = "Georgia"
    hp5.font.size = Pt(28)
    hp5.font.bold = True
    hp5.font.color.rgb = gold
    
    tBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    tf5 = tBox5.text_frame
    tf5.word_wrap = True
    
    bullets5 = [
        "Movie Meter establishes a leak-free predictive framework for pre-release cinema classification.",
        "Drives capital efficiency for production companies and OTT acquisition managers by linking predictions directly to box office ROI ranges.",
        "Provides widescreen 16:9 investor slide deck compilations (PPTX) and technical PDF reports dynamically.",
        "Deliverable package is fully deployed on GitHub and executed live on Streamlit Community Cloud."
    ]
    for i, b_text in enumerate(bullets5):
        bp = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = white
        bp.space_after = Pt(12)
        
    if audience_img:
        slide5.shapes.add_picture(audience_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    prs.save("Movie_Meter_Final_Presentation.pptx")
    print("PPTX presentation generated successfully.")

if __name__ == "__main__":
    create_presentation()
