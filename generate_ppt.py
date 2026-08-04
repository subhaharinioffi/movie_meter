import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Aspect Ratio
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    maroon = RGBColor(109, 0, 26)     # #6D001A
    gold = RGBColor(255, 213, 79)    # #FFD54F
    slate = RGBColor(71, 85, 105)     # #475569
    dark_blue = RGBColor(15, 23, 42)  # #0F172A
    white = RGBColor(255, 255, 255)
    off_white = RGBColor(250, 250, 250)
    
    # Locate generated images in brain directory
    brain_dir = r"C:\Users\Subhaharini\.gemini\antigravity\brain\3651ea6e-94c7-4edf-89e9-b55af2040db7"
    
    def get_image(prefix):
        files = glob.glob(os.path.join(brain_dir, f"{prefix}_*.jpg"))
        return files[0] if files else None
        
    hero_img = get_image("hero_banner_mockup")
    showcase_img = get_image("movie_showcase_mockup")
    gauge_img = get_image("success_gauge_mockup")
    revenue_img = get_image("revenue_chart_mockup")
    audience_img = get_image("audience_insights_mockup")
    
    # ------------------ SLIDE 1: Title Slide (Dark Maroon Theme) ------------------
    slide_layout = prs.slide_layouts[6] # Blank Slide
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = maroon
    
    # Title Text Box
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MOVIE METER"
    p.font.name = "Georgia"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = gold
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Cinema Success\nAnalytics Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = white
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Designed for South Indian Cinema Catalog Acquisitions"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = off_white
    
    # Add Image Mockup
    if hero_img:
        slide1.shapes.add_picture(hero_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 2: Business Problem (Off-white Theme) ------------------
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = off_white
    
    # Heading
    hBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf2 = hBox2.text_frame
    hp2 = htf2.paragraphs[0]
    hp2.text = "The Business Problem: Capital Risk in Acquisitions"
    hp2.font.name = "Georgia"
    hp2.font.size = Pt(28)
    hp2.font.bold = True
    hp2.font.color.rgb = maroon
    
    # Bullet points
    tBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.5))
    tf2 = tBox2.text_frame
    tf2.word_wrap = True
    
    bullets2 = [
        "Traditional catalog licensing for OTT platforms carries severe premium cost risks.",
        "Pre-production budget planning is frequently based on subjective actor/director packaging.",
        "Standard predictive systems rely on ticket counts or ratings, which suffer from data leakage.",
        "Movie Meter provides a standardized benchmark framework for movie quality categorization before release."
    ]
    for i, b_text in enumerate(bullets2):
        bp = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(16)
        bp.font.color.rgb = slate
        bp.space_after = Pt(14)
        
    # Add Image Mockup
    if showcase_img:
        slide2.shapes.add_picture(showcase_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 3: ML Pipeline & Dataflow (Off-white Theme) ------------------
    slide3 = prs.slides.add_slide(slide_layout)
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = off_white
    
    hBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf3 = hBox3.text_frame
    hp3 = htf3.paragraphs[0]
    hp3.text = "Technical Architecture & Data Pipeline"
    hp3.font.name = "Georgia"
    hp3.font.size = Pt(28)
    hp3.font.bold = True
    hp3.font.color.rgb = maroon
    
    tBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
    tf3 = tBox3.text_frame
    tf3.word_wrap = True
    
    bullets3 = [
        "Memory-Efficient IMDb TSV Datasets Joining: Extracts region and language codes across 1.3 GB of official records chunk-by-chunk under low RAM bounds.",
        "smoothed Out-of-Fold Target Encoding: Calculates cross-validated rating distributions for actors and directors to index their reputation without data leakage.",
        "Custom preprocessing: Cleans duration profiles using median estimators and bins content ratings into certification nodes (G, PG, PG-13, R)."
    ]
    for i, b_text in enumerate(bullets3):
        bp = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(16)
        bp.font.color.rgb = slate
        bp.space_after = Pt(16)
        
    # Simple ASCII Diagram box
    flowBox = slide3.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0))
    ftf = flowBox.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "IMDb Basics & Ratings  ──>  Filter (Tamil/IN releases)  ──>  smoothed Encoding  ──>  XGBoost Classifier"
    fp.font.name = "Courier New"
    fp.font.size = Pt(14)
    fp.font.bold = True
    fp.font.color.rgb = maroon
    fp.alignment = PP_ALIGN.CENTER
    
    # ------------------ SLIDE 4: Machine Learning Model (Off-white Theme) ------------------
    slide4 = prs.slides.add_slide(slide_layout)
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = off_white
    
    hBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf4 = hBox4.text_frame
    hp4 = htf4.paragraphs[0]
    hp4.text = "Model Benchmarks & Hyperparameters"
    hp4.font.name = "Georgia"
    hp4.font.size = Pt(28)
    hp4.font.bold = True
    hp4.font.color.rgb = maroon
    
    tBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.5))
    tf4 = tBox4.text_frame
    tf4.word_wrap = True
    
    bullets4 = [
        "XGBoost Classifier selected for classification boundary accuracy.",
        "5-Fold Stratified CV Accuracy: 62.1% | ROC-AUC: 73.0%.",
        "Outperforms baseline Logistic Regression configurations by ~12% on low/high class margins.",
        "Optimal parameters: estimators=300, max_depth=6, learning_rate=0.03, subsample=0.8."
    ]
    for i, b_text in enumerate(bullets4):
        bp = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(16)
        bp.font.color.rgb = slate
        bp.space_after = Pt(14)
        
    # Add Image Mockup
    if gauge_img:
        slide4.shapes.add_picture(gauge_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    # ------------------ SLIDE 5: Tech Stack & Visual Interface (Off-white Theme) ------------------
    slide5 = prs.slides.add_slide(slide_layout)
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = off_white
    
    hBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    htf5 = hBox5.text_frame
    hp5 = htf5.paragraphs[0]
    hp5.text = "The Tech Stack & Enterprise SaaS UI"
    hp5.font.name = "Georgia"
    hp5.font.size = Pt(28)
    hp5.font.bold = True
    hp5.font.color.rgb = maroon
    
    tBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(4.5))
    tf5 = tBox5.text_frame
    tf5.word_wrap = True
    
    bullets5 = [
        "Frontend Interface: Built using Streamlit framework with responsive HTML container overrides.",
        "SaaS Aesthetic: Inter and Poppins font typography paired with Font Awesome v6 vector icons.",
        "Interactive Plots: Plotly Express gauges and histograms.",
        "Digital Match: Automated OTT acquisition recommendation index logic."
    ]
    for i, b_text in enumerate(bullets5):
        bp = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        bp.text = f"• {b_text}"
        bp.font.name = "Arial"
        bp.font.size = Pt(16)
        bp.font.color.rgb = slate
        bp.space_after = Pt(14)
        
    # Add Image Mockup
    if revenue_img:
        slide5.shapes.add_picture(revenue_img, Inches(7.5), Inches(1.8), width=Inches(5.0))
        
    prs.save("Movie_Meter_Presentation.pptx")
    print("PPTX presentation generated successfully.")

if __name__ == "__main__":
    create_presentation()
